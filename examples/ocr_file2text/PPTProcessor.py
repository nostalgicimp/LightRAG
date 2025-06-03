import os
import io
from typing import Tuple, Optional, Dict, Any, List
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE # 用于识别形状类型，如图片
from pptx.exc import PackageNotFoundError  # python-pptx 特定的异常

from PIL import Image # Pillow for image manipulation
import pytesseract     # Tesseract OCR
import numpy as np     # PaddleOCR needs numpy

# 尝试导入PaddleOCR并设置标志
try:
    from paddleocr import PaddleOCR
    PADDLEOCR_AVAILABLE = True
except ImportError:
    PADDLEOCR_AVAILABLE = False
    PaddleOCR = None # 定义PaddleOCR为None，以避免后续代码出现NameError

class PPTProcessor:
    """
    处理 PowerPoint (.pptx) 文件内容，提取文本数据，包括幻灯片上的文本和图像中的文本 (OCR)。
    对 .ppt 文件提供转换建议。
    """

    def __init__(self,
                 tesseract_cmd_path: Optional[str] = None,
                 default_ocr_engine: str = 'paddle', # 默认OCR引擎
                 paddle_ocr_config: Optional[Dict[str, Any]] = None):
        """
        初始化 PPTProcessor。

        Args:
            tesseract_cmd_path (Optional[str]): Tesseract OCR 的可执行文件路径。
                                                如果 Tesseract 已在系统 PATH 中，则无需设置。
            default_ocr_engine (str): 默认使用的OCR引擎 ('tesseract' 或 'paddle')。
            paddle_ocr_config (Optional[Dict[str, Any]]): PaddleOCR 的初始化配置。
                例如: {'use_angle_cls': True, 'lang': 'ch', 'show_log': False}
        """
        self.tesseract_available = False
        if tesseract_cmd_path:
            pytesseract.pytesseract.tesseract_cmd = tesseract_cmd_path

        try:
            pytesseract.get_tesseract_version()
            self.tesseract_available = True
        except pytesseract.TesseractNotFoundError:
            print("警告: Tesseract OCR 未找到或未正确配置。如果选择Tesseract，图像文本提取功能将不可用。")
        except Exception as e:
            print(f"警告: 初始化Tesseract时发生错误: {e}。Tesseract图像文本提取功能可能受影响。")

        self.default_ocr_engine = default_ocr_engine.lower()
        if self.default_ocr_engine == 'paddle' and not PADDLEOCR_AVAILABLE:
            print("警告: PaddleOCR库未安装，但被选为默认OCR引擎。"
                  "将尝试回退到Tesseract (如果可用)，否则不进行图像OCR。")
            if self.tesseract_available:
                self.default_ocr_engine = 'tesseract'
            else:
                self.default_ocr_engine = None

        self.paddle_ocr_instance: Optional[PaddleOCR] = None
        self.user_paddle_ocr_config = paddle_ocr_config
        self.current_paddle_lang_init = None

    def _initialize_paddleocr(self, lang_for_paddle: str) -> bool:
        """【私有方法】惰性初始化PaddleOCR实例。"""
        if not PADDLEOCR_AVAILABLE:
            print("警告: PaddleOCR库未安装，无法使用PaddleOCR引擎。")
            return False

        if self.user_paddle_ocr_config:
            should_reinitialize_user_config = self.paddle_ocr_instance is None
            if 'lang' in self.user_paddle_ocr_config:
                 should_reinitialize_user_config = should_reinitialize_user_config or \
                    (self.current_paddle_lang_init != self.user_paddle_ocr_config.get('lang'))
            if should_reinitialize_user_config:
                try:
                    print(f"正在使用用户配置初始化PaddleOCR: {self.user_paddle_ocr_config}...")
                    self.paddle_ocr_instance = PaddleOCR(**self.user_paddle_ocr_config)
                    self.current_paddle_lang_init = self.user_paddle_ocr_config.get('lang', 'ch')
                    print("PaddleOCR已通过用户配置成功初始化。")
                    return True
                except Exception as e:
                    print(f"警告: 使用用户配置初始化PaddleOCR失败: {e}")
                    self.paddle_ocr_instance = None
                    return False
            return True

        if self.paddle_ocr_instance is None or self.current_paddle_lang_init != lang_for_paddle:
            try:
                default_config = {'use_angle_cls': True, 'lang': lang_for_paddle, 'show_log': False}
                print(f"正在使用默认配置初始化PaddleOCR (语言: {lang_for_paddle})...")
                self.paddle_ocr_instance = PaddleOCR(**default_config)
                self.current_paddle_lang_init = lang_for_paddle
                print("PaddleOCR已成功初始化。")
                return True
            except Exception as e:
                print(f"警告: 初始化PaddleOCR失败 (语言={lang_for_paddle}): {e}")
                self.paddle_ocr_instance = None
                return False
        return True

    def _map_tesseract_lang_to_paddle(self, tesseract_langs: str) -> str:
        """【私有方法】启发式地将Tesseract语言字符串映射到PaddleOCR语言代码。"""
        if self.user_paddle_ocr_config and 'lang' in self.user_paddle_ocr_config:
            return self.user_paddle_ocr_config['lang']
        t_langs = tesseract_langs.lower()
        if 'chi_sim' in t_langs or 'chi_tra' in t_langs: return 'ch'
        if 'eng' in t_langs: return 'en'
        if 'kor' in t_langs: return 'korean'
        if 'japan' in t_langs: return 'japan'
        print(f"警告: 无法将Tesseract语言 '{tesseract_langs}' 直接映射到PaddleOCR支持的语言。"
              f"将默认使用 'ch'。")
        return 'ch'

    def _ocr_image_with_tesseract_from_bytes(self, image_bytes: bytes, languages: str, image_id: str) -> str:
        """【私有方法】使用Tesseract从图像字节流OCR。"""
        if not self.tesseract_available:
            print(f"  - Tesseract OCR 不可用，无法处理图像 (ID: {image_id})。")
            return ""
        try:
            pil_image = Image.open(io.BytesIO(image_bytes))
            if pil_image.mode == 'RGBA' or pil_image.mode == 'P':
                pil_image = pil_image.convert('RGB')
            text = pytesseract.image_to_string(pil_image, lang=languages)
            return text.strip()
        except pytesseract.TesseractError as te:
            print(f"  - Tesseract OCR错误 (图像ID: {image_id}): {te}")
        except Image.UnidentifiedImageError:
            print(f"  - Tesseract警告: 无法识别图像格式 (图像ID: {image_id})。")
        except Exception as e_img:
            print(f"  - Tesseract处理图像 (图像ID: {image_id}) 时发生未知错误: {e_img}")
        return ""

    def _ocr_image_with_paddle_from_bytes(self, image_bytes: bytes, lang_for_paddle: str, image_id: str) -> str:
        """【私有方法】使用PaddleOCR从图像字节流OCR。"""
        if not self._initialize_paddleocr(lang_for_paddle) or not self.paddle_ocr_instance:
            print(f"  - PaddleOCR初始化失败或不可用，无法处理图像 (图像ID: {image_id})")
            return ""
        try:
            pil_image = Image.open(io.BytesIO(image_bytes))
            if pil_image.mode == 'RGBA' or pil_image.mode == 'P':
                pil_image = pil_image.convert('RGB')
            elif pil_image.mode == 'L': # 灰度图
                 pil_image = pil_image.convert('RGB')

            img_np = np.array(pil_image)
            ocr_results = self.paddle_ocr_instance.ocr(img_np, cls=True)

            image_texts = []
            if ocr_results and ocr_results[0] is not None:
                for line_info in ocr_results[0]:
                    text_content = line_info[1][0]
                    if text_content.strip():
                        image_texts.append(text_content.strip())
            return "\n".join(image_texts)
        except Image.UnidentifiedImageError:
            print(f"  - PaddleOCR警告: 无法识别图像格式 (图像ID: {image_id})。")
        except Exception as e_img:
            print(f"  - PaddleOCR处理图像 (图像ID: {image_id}) 时发生错误: {e_img}")
        return ""

    def extract_text_from_ppt(self,
                              ppt_path: str,
                              include_metadata: bool = False,
                              include_notes: bool = True,
                              include_images_ocr: bool = True,
                              ocr_languages: str = 'chi_sim+eng',
                              ocr_engine: Optional[str] = None
                             ) -> Tuple[str, Optional[Dict[str, Any]]]:
        """
        从 PowerPoint (.pptx) 文件中提取文本内容，包括OCR图像文本。
        对于 .ppt 文件，会给出转换建议。

        参数:
            ppt_path (str): .pptx (或 .ppt) 文件路径。
            include_metadata (bool): 是否包含文件元数据。
            include_notes (bool): 是否包含幻灯片备注中的文本。
            include_images_ocr (bool): 是否尝试提取图像中的文本。
            ocr_languages (str): OCR识别时使用的语言。
            ocr_engine (Optional[str]): 使用的OCR引擎 ('tesseract', 'paddle')。如果None，使用默认引擎。

        返回:
            元组: (提取的文本内容, 元数据字典或None)
        """
        file_basename = os.path.basename(ppt_path)
        if ppt_path.lower().endswith(".ppt"):
            # python-pptx 不支持旧的 .ppt 格式。
            # 可以考虑的方案：
            # 1. 提示用户手动转换。
            # 2. (高级) 集成 unoconv (Linux/macOS) 或 pywin32 (Windows, 需安装Office) 来自动转换。
            #    这会增加外部依赖和复杂性。例如，使用unoconv:
            #    import subprocess, shutil, tempfile
            #    if shutil.which("unoconv"):
            #        # ... 转换逻辑 ...
            #        pass # 然后递归调用本方法处理转换后的 .pptx
            #    else:
            #        raise ValueError(...)
            error_msg = (
                f"错误：检测到旧版 .ppt 文件 ({file_basename})。\n"
                "此处理器基于 python-pptx，仅直接支持现代 .pptx 格式。\n"
                "请先将 .ppt 文件转换为 .pptx 格式（例如，使用PowerPoint或LibreOffice另存为）。"
            )
            raise ValueError(error_msg)
        elif not ppt_path.lower().endswith(".pptx"):
            raise ValueError(f"错误：不支持的文件类型 '{file_basename}'。此处理器仅支持 .pptx 文件。")

        all_text_parts: List[str] = []
        metadata_dict: Optional[Dict[str, Any]] = None
        chosen_ocr_engine = (ocr_engine or self.default_ocr_engine or "none").lower()

        try:
            print(f"正在处理PPTX文件: {file_basename}...")
            prs = Presentation(ppt_path)

            if include_metadata:
                metadata_dict = self._get_ppt_metadata(prs, ppt_path)

            for slide_idx, slide in enumerate(prs.slides):
                slide_text_collector: List[str] = []
                image_count_on_slide = 0 # 用于为图像生成唯一ID

                # 提取幻灯片上各个形状中的文本
                for shape_idx, shape in enumerate(slide.shapes):
                    if shape.has_text_frame:
                        text_in_shape = shape.text_frame.text.strip()
                        if text_in_shape:
                            slide_text_collector.append(text_in_shape)

                    # 处理图片 (如果形状是图片类型且需要OCR)
                    if include_images_ocr and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_count_on_slide += 1
                        image_id_for_log = f"幻灯片{slide_idx + 1}-图片{image_count_on_slide}"
                        
                        ocr_text_from_image = ""
                        try:
                            image = shape.image
                            image_bytes = image.blob # 图像的字节数据
                            # image_format_ext = image.ext # 例如 'png', 'jpg'

                            if chosen_ocr_engine == "none":
                                # print(f"信息: 未选择OCR引擎，跳过图像 {image_id_for_log} 的文本提取。")
                                pass
                            elif chosen_ocr_engine == 'tesseract':
                                ocr_text_from_image = self._ocr_image_with_tesseract_from_bytes(
                                    image_bytes, ocr_languages, image_id_for_log
                                )
                            elif chosen_ocr_engine == 'paddle':
                                lang_for_paddle = self._map_tesseract_lang_to_paddle(ocr_languages)
                                ocr_text_from_image = self._ocr_image_with_paddle_from_bytes(
                                    image_bytes, lang_for_paddle, image_id_for_log
                                )
                            else:
                                if chosen_ocr_engine not in [None, "none"]: # 避免重复打印
                                     print(f"警告: 未知的OCR引擎 '{chosen_ocr_engine}'。跳过图像 {image_id_for_log} 的文本提取。")
                        
                        except AttributeError: # 有些“图片”可能没有 .image.blob (例如，链接的图片或某些特殊图形)
                            print(f"警告: 无法获取图像数据 (ID: {image_id_for_log})。可能不是嵌入式图片。")
                        except Exception as e_img_extract:
                            print(f"警告: 提取或处理图像 {image_id_for_log} 时发生错误: {e_img_extract}")

                        if ocr_text_from_image:
                            slide_text_collector.append(f"\n--- 来自图像 {image_id_for_log} (OCR: {chosen_ocr_engine}) ---\n{ocr_text_from_image.strip()}")
                
                if slide_text_collector: # 如果当前幻灯片收集到了文本（包括OCR的）
                    all_text_parts.append(f"\n--- 幻灯片 {slide_idx + 1} ---\n" + "\n\n".join(filter(None, slide_text_collector)))

                # 提取备注
                if include_notes and slide.has_notes_slide:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame and notes_slide.notes_text_frame.text:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            all_text_parts.append(f"\n--- 幻灯片 {slide_idx + 1} 的备注 ---\n{notes_text}")
            
            full_text = "\n".join(filter(None, all_text_parts)).strip()
            print(f"文件 {file_basename} 处理完成。")

        except PackageNotFoundError:
            raise FileNotFoundError(f"错误：PPTX文件未找到或不是有效的 .pptx 格式: {ppt_path}")
        except ValueError as ve: # 捕获我们自己为 .ppt 文件抛出的错误
            raise ve
        except Exception as e:
            raise Exception(f"处理PPT文件 '{ppt_path}' 时发生一般性错误: {str(e)}")

        return full_text, metadata_dict

    def _get_ppt_metadata(self, prs: Presentation, ppt_path: str) -> Dict[str, Any]:
        """获取 .pptx 文件的元数据。"""
        file_stats = os.stat(ppt_path)
        metadata = {
            "文件名称": os.path.basename(ppt_path),
            "文件大小_字节": file_stats.st_size,
            "创建时间_本地": file_stats.st_ctime,
            "修改时间_本地": file_stats.st_mtime,
            "幻灯片数量": len(prs.slides)
        }
        cp = prs.core_properties
        core_prop_map = {
            "作者": cp.author, "类别": cp.category, "备注": cp.comments,
            "内容状态": cp.content_status,
            "创建日期_文档": cp.created.isoformat() if cp.created else None,
            "标识符": cp.identifier, "关键词": cp.keywords, "语言": cp.language,
            "最后修改者": cp.last_modified_by,
            "最后打印日期": cp.last_printed.isoformat() if cp.last_printed else None,
            "修改日期_文档": cp.modified.isoformat() if cp.modified else None,
            "修订号": cp.revision, "主题": cp.subject, "标题": cp.title, "版本": cp.version,
        }
        metadata.update({k: v for k, v in core_prop_map.items() if v}) # 只添加有值的元数据
        return metadata


